from pptx import Presentation
from pptx import Presentation
from pptx.util import Pt

import six
import copy
import sys
import getopt
import os.path


def _get_blank_slide_layout(pres):
	layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
	min_items = min(layout_items_count)
	blank_layout_id = layout_items_count.index(min_items)
	return pres.slide_layouts[blank_layout_id]

def duplicate_slide(pres, index):
	"""Duplicate the slide with the given index in pres.

	Adds slide to the end of the presentation"""
	source = pres.slides[index]

	blank_slide_layout = _get_blank_slide_layout(pres)
	dest = pres.slides.add_slide(blank_slide_layout)

	for shp in source.shapes:
		el = shp.element
		newel = copy.deepcopy(el)
		dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

	for key, value in six.iteritems(source.part.rels):
		# Make sure we don't copy a notesSlide relation as that won't exist
		if not "notesSlide" in value.reltype:
			dest.part.rels.add_relationship(value.reltype, value._target, value.rId)

	return dest

def main():
	xlsx_input = ''
	pptx_input = '/../Resultatlista.pptx'

	prs = Presentation(os.path.dirname(__file__) + pptx_input)

	duplicate_slide(prs, 0)

	prs.save(os.path.dirname(__file__) + pptx_input)

	

main()